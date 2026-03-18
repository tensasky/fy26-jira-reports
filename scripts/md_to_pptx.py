#!/usr/bin/env python3
"""
将 Markdown PPT 转换为 PowerPoint (.pptx)
"""
import sys
import re
from pathlib import Path

def parse_markdown(md_content):
    """解析 Markdown 内容为幻灯片结构"""
    slides = []
    current_slide = None
    
    lines = md_content.split('\n')
    
    for line in lines:
        line = line.strip()
        
        # 新的幻灯片标题 (## 幻灯片 X: 标题)
        if line.startswith('## 幻灯片'):
            if current_slide:
                slides.append(current_slide)
            title = line.split(':', 1)[1].strip() if ':' in line else line
            current_slide = {'title': title, 'content': []}
        
        # 幻灯片标题 (heading 级别)
        elif line.startswith('# ') and '幻灯片' not in line:
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': line[2:].strip(), 'content': []}
        
        # 副标题
        elif line.startswith('## ') and current_slide:
            current_slide['content'].append(('heading', line[2:].strip()))
        
        # 列表项
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide:
                current_slide['content'].append(('bullet', line[2:].strip()))
        
        # 表格行
        elif line.startswith('|') and current_slide:
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if cells and not all(c in ['---', ''] for c in cells):
                current_slide['content'].append(('table_row', cells))
        
        # 普通段落
        elif line and current_slide:
            # 跳过代码块标记
            if line.startswith('```'):
                continue
            current_slide['content'].append(('text', line))
    
    if current_slide:
        slides.append(current_slide)
    
    return slides

def create_pptx(slides, output_path):
    """创建 PowerPoint 文件"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        print("❌ 需要安装 python-pptx 库")
        print("请运行: pip3 install python-pptx")
        return False
    
    prs = Presentation()
    
    # 设置幻灯片尺寸为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        # 添加空白幻灯片
        blank_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # 添加内容
        if slide_data['content']:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.7)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for item_type, item_text in slide_data['content']:
                p = content_frame.add_paragraph()
                
                if item_type == 'heading':
                    p.text = item_text
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 102, 153)
                    p.space_before = Pt(12)
                
                elif item_type == 'bullet':
                    p.text = "• " + item_text
                    p.font.size = Pt(18)
                    p.level = 0
                    p.space_before = Pt(6)
                
                elif item_type == 'table_row':
                    p.text = ' | '.join(item_text)
                    p.font.size = Pt(14)
                    p.font.italic = True
                    p.space_before = Pt(4)
                
                else:  # text
                    p.text = str(item_text)
                    p.font.size = Pt(16)
                    p.space_before = Pt(4)
    
    prs.save(output_path)
    return True

def convert_md_to_pptx(md_file, pptx_file=None):
    """转换单个 Markdown 文件到 PPTX"""
    md_path = Path(md_file)
    
    if not md_path.exists():
        print(f"❌ 文件不存在: {md_file}")
        return False
    
    if pptx_file is None:
        pptx_file = md_path.with_suffix('.pptx')
    
    print(f"📄 读取: {md_file}")
    md_content = md_path.read_text(encoding='utf-8')
    
    print(f"🔍 解析内容...")
    slides = parse_markdown(md_content)
    print(f"✅ 解析完成: {len(slides)} 张幻灯片")
    
    print(f"💾 生成: {pptx_file}")
    if create_pptx(slides, pptx_file):
        print(f"✅ 完成: {pptx_file}")
        return True
    else:
        print("❌ 生成失败")
        return False

def main():
    """主函数"""
    output_dir = Path('/Users/admin/.openclaw/workspace/output')
    
    md_files = [
        'ppt_sparkle_leadership.md',
        'ppt_onboarding_teams.md',
        'ppt_standards_admins.md',
        'ppt_knowledge_base_users.md'
    ]
    
    success_count = 0
    
    for md_file in md_files:
        md_path = output_dir / md_file
        if md_path.exists():
            print(f"\n{'='*50}")
            if convert_md_to_pptx(md_path):
                success_count += 1
        else:
            print(f"❌ 跳过: {md_file} (不存在)")
    
    print(f"\n{'='*50}")
    print(f"🎉 完成: {success_count}/{len(md_files)} 个文件转换成功")
    
    return success_count

if __name__ == '__main__':
    main()
